#!/usr/bin/env python3
"""
Pitch Deck Template Generator for Infinite Voyage

Generates professional pitch presentation templates (.pptx) for publisher pitches,
investor presentations, and team presentations. Each deck includes branded slides
with speaker notes, consistent styling, and appropriate structure for its audience.

Usage:
    python generate_templates.py
    python generate_templates.py --output-dir ./output

Requirements:
    pip install python-pptx
"""

import argparse
import os
import sys
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(
        "Error: python-pptx is not installed.\n"
        "Install it with:\n"
        "    pip install python-pptx\n"
        "\n"
        "Or if using a virtual environment:\n"
        "    python -m pip install python-pptx"
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# Shared styling constants
# ---------------------------------------------------------------------------

BRAND_PRIMARY = RGBColor(0, 102, 255)
BRAND_SECONDARY = RGBColor(255, 153, 0)
BRAND_DARK = RGBColor(20, 20, 40)
BRAND_LIGHT = RGBColor(240, 245, 255)
BRAND_WHITE = RGBColor(255, 255, 255)
BRAND_GRAY = RGBColor(150, 150, 150)
BRAND_GREEN = RGBColor(46, 139, 87)
BRAND_RED = RGBColor(200, 50, 50)

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

TODAY_DISPLAY = datetime.now().strftime("%B %Y")
TODAY_ISO = datetime.now().strftime("%Y-%m-%d")


def new_presentation():
    """Create a widescreen 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_background(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BRAND_DARK, alignment=PP_ALIGN.LEFT,
                font_name=None):
    """Add a formatted text box to a slide."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return txbox


def add_bullet_slide(prs, title, bullets, speaker_notes="",
                     bg_color=BRAND_WHITE, title_color=BRAND_DARK,
                     bullet_color=BRAND_DARK):
    """Add a standard title + bullet-list slide with speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_color)

    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, title,
                font_size=36, bold=True, color=title_color)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.05), Inches(8.8), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_PRIMARY
    line.line.fill.background()

    y = 1.3
    for bullet_text in bullets:
        add_textbox(slide, 0.9, y, 8.2, 0.4, f"  {bullet_text}",
                    font_size=18, color=bullet_color)
        y += 0.45

    if speaker_notes:
        notes = slide.notes_slide
        notes.notes_text_frame.text = speaker_notes

    return slide


def add_title_slide(prs, title, subtitle, extra_line="",
                    bg_color=BRAND_DARK, title_color=BRAND_PRIMARY,
                    subtitle_color=BRAND_WHITE, notes=""):
    """Add a branded title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_color)

    add_textbox(slide, 0.5, 1.2, 9.0, 1.2, title,
                font_size=54, bold=True, color=title_color,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 2.5, 9.0, 0.6, subtitle,
                font_size=24, color=subtitle_color,
                alignment=PP_ALIGN.CENTER)

    if extra_line:
        add_textbox(slide, 0.5, 4.6, 9.0, 0.5, extra_line,
                    font_size=14, color=BRAND_GRAY,
                    alignment=PP_ALIGN.CENTER)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_thank_you_slide(prs, contact_name="[Your Name]",
                        contact_email="[email@studio.com]",
                        contact_url="[infinitevoyage.com]"):
    """Add a closing thank-you slide with contact info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_DARK)

    add_textbox(slide, 0.5, 1.2, 9.0, 1.0, "Thank You",
                font_size=54, bold=True, color=BRAND_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    contact = f"Questions?\n\n{contact_name}\n{contact_email}\n{contact_url}"
    txbox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.8), Inches(7.0), Inches(2.0)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = contact
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.color.rgb = BRAND_WHITE
            run.font.size = Pt(16)

    slide.notes_slide.notes_text_frame.text = (
        "Thank the audience. Open the floor for Q&A. "
        "Have backup slides ready for detailed questions."
    )
    return slide


# ===========================================================================
# Template generators
# ===========================================================================


def generate_publisher_pitch(output_dir):
    """Generate publisher-pitch-template.pptx."""
    prs = new_presentation()

    add_title_slide(prs, "INFINITE VOYAGE", "Publisher Pitch Deck",
                    f"[Studio Name] | {TODAY_DISPLAY}")

    add_bullet_slide(prs, "Core Gameplay", [
        "Genre: [e.g., Narrative Space Exploration RPG]",
        "Core Loop: Explore > Discover > Research > Upgrade",
        "Key Mechanic: [Your unique selling mechanic]",
        "Session Length: [e.g., 30-60 min]",
        "Comparable Titles: [e.g., Outer Wilds meets Mass Effect]",
    ])

    add_bullet_slide(prs, "Market Opportunity", [
        "Target Market: [Market segment and size]",
        "Growth Rate: [YoY growth]",
        "Comparable Sales: [Reference title sales figures]",
    ])

    add_bullet_slide(prs, "The Team", [
        "[Lead Name] - Game Director",
        "[Lead Name] - Lead Programmer",
        "[Lead Name] - Art Director",
        "Team Size: [X core + X contractors]",
    ])

    add_bullet_slide(prs, "Timeline & Ask", [
        "Current Status: [Prototype / Alpha / Beta]",
        "Launch Target: [Date]",
        "Funding Need: [$X M]",
        "Partnership Type: [Publishing Deal]",
    ])

    add_thank_you_slide(prs)

    filepath = os.path.join(output_dir, "publisher-pitch-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


def generate_investor_pitch(output_dir):
    """Generate investor-pitch-template.pptx."""
    prs = new_presentation()

    add_title_slide(prs, "INFINITE VOYAGE", "Investor Presentation",
                    f"[Studio Name] | {TODAY_DISPLAY} | Confidential")

    add_bullet_slide(prs, "The Problem", [
        "[Market gap]",
        "[Unmet player demand]",
        "[Why now is the right time]",
    ])

    add_bullet_slide(prs, "Our Solution", [
        "[Unique value proposition]",
        "[Core gameplay experience]",
        "[Why players will choose this]",
    ])

    add_bullet_slide(prs, "Market Size", [
        "TAM: [$XX B]",
        "SAM: [$XX B]",
        "SOM: [$XX M]",
    ])

    add_bullet_slide(prs, "Traction", [
        "Playtest results",
        "Wishlists / community",
        "Awards / press",
    ])

    add_bullet_slide(prs, "Financial Projections", [
        "Year 1 Revenue: [$X M]",
        "Year 2 Revenue: [$X M]",
        "Break-even: [Month X]",
    ])

    add_bullet_slide(prs, "The Ask", [
        "Funding: [$X.X M]",
        "Use of Funds: Development X%, Marketing X%, Operations X%",
        "Terms: [Equity / Revenue Share]",
    ])

    add_thank_you_slide(prs)

    filepath = os.path.join(output_dir, "investor-pitch-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


def generate_team_presentation(output_dir):
    """Generate team-presentation-template.pptx."""
    prs = new_presentation()

    add_title_slide(prs, "INFINITE VOYAGE", "Team Presentation",
                    f"[Studio Name] | {TODAY_DISPLAY} | Internal")

    add_bullet_slide(prs, "Vision & Pillars", [
        "Pillar 1: [e.g., Exploration]",
        "Pillar 2: [e.g., Discovery]",
        "Pillar 3: [e.g., Consequence]",
    ])

    add_bullet_slide(prs, "Core Gameplay Loop", [
        "EXPLORE: Navigate star systems",
        "DISCOVER: Investigate ruins and artifacts",
        "RESEARCH: Unlock technology",
        "UPGRADE: Improve ship and abilities",
    ])

    add_bullet_slide(prs, "Systems Overview", [
        "Combat, Crafting, Skill Trees",
        "Equipment, Economy, AI",
    ])

    add_bullet_slide(prs, "Roadmap", [
        "Milestone 1 (Done): Pre-production",
        "Milestone 2 (Current): Vertical slice",
        "Milestone 3: Alpha",
        "Milestone 4: Beta",
        "Milestone 5: Launch",
    ])

    add_bullet_slide(prs, "Action Items", [
        "[Action 1]: [Owner] -- Due [Date]",
        "[Action 2]: [Owner] -- Due [Date]",
        "[Action 3]: [Owner] -- Due [Date]",
    ])

    add_thank_you_slide(prs)

    filepath = os.path.join(output_dir, "team-presentation-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


# ===========================================================================
# CLI entry point
# ===========================================================================


def main():
    parser = argparse.ArgumentParser(
        description="Generate Infinite Voyage pitch deck presentation templates (.pptx).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python generate_templates.py\n"
            "  python generate_templates.py --output-dir ./output\n"
        ),
    )
    parser.add_argument(
        "--output-dir",
        default=os.path.dirname(os.path.abspath(__file__)),
        help="Directory where .pptx files will be written (default: script directory).",
    )
    args = parser.parse_args()

    output_dir = os.path.abspath(args.output_dir)
    os.makedirs(output_dir, exist_ok=True)

    print(f"Generating pitch deck templates in: {output_dir}\n")

    generate_publisher_pitch(output_dir)
    generate_investor_pitch(output_dir)
    generate_team_presentation(output_dir)

    print("\nDone. All pitch deck templates generated successfully.")


if __name__ == "__main__":
    main()
