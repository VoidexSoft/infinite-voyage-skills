#!/usr/bin/env python3
"""
Pitch Deck Template Generator for Infinite Voyage

Generates professional pitch presentation templates (.pptx) for publisher pitches,
investor presentations, and team presentations. Each deck includes branded slides
with speaker notes, consistent styling, and appropriate structure for its audience.

Usage:
    python generate_templates.py
    python generate_templates.py --output-dir ./output

Requirements:
    pip install python-pptx
"""

import argparse
import os
import sys
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(
        "Error: python-pptx is not installed.\n"
        "Install it with:\n"
        "    pip install python-pptx\n"
        "\n"
        "Or if using a virtual environment:\n"
        "    python -m pip install python-pptx"
    )
    sys.exit(1)


# ---------------------------------------------------------------------------
# Shared styling constants
# ---------------------------------------------------------------------------

BRAND_PRIMARY = RGBColor(0, 102, 255)   # Blue
BRAND_SECONDARY = RGBColor(255, 153, 0) # Orange
BRAND_DARK = RGBColor(20, 20, 40)       # Near-black
BRAND_LIGHT = RGBColor(240, 245, 255)   # Light blue
BRAND_WHITE = RGBColor(255, 255, 255)
BRAND_GRAY = RGBColor(150, 150, 150)
BRAND_GREEN = RGBColor(46, 139, 87)
BRAND_RED = RGBColor(200, 50, 50)

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)  # 16:9

TODAY_DISPLAY = datetime.now().strftime("%B %Y")
TODAY_ISO = datetime.now().strftime("%Y-%m-%d")


def new_presentation():
    """Create a widescreen 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_background(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BRAND_DARK, alignment=PP_ALIGN.LEFT,
                font_name=None):
    """Add a formatted text box to a slide."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = text
    para = tf.paragraphs[0]
    para.alignment = alignment
    run = para.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name
    return txbox


def add_bullet_slide(prs, title, bullets, speaker_notes="",
                     bg_color=BRAND_WHITE, title_color=BRAND_DARK,
                     bullet_color=BRAND_DARK):
    """Add a standard title + bullet-list slide with speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, bg_color)

    # Title
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, title,
                font_size=36, bold=True, color=title_color)

    # Accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.05), Inches(8.8), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_PRIMARY
    line.line.fill.background()

    # Bullets
    y = 1.3
    for bullet_text in bullets:
        add_textbox(slide, 0.9, y, 8.2, 0.4, f"  {bullet_text}",
                    font_size=18, color=bullet_color)
        y += 0.45

    # Speaker notes
    if speaker_notes:
        notes = slide.notes_slide
        notes.notes_text_frame.text = speaker_notes

    return slide


def add_title_slide(prs, title, subtitle, extra_line="",
                    bg_color=BRAND_DARK, title_color=BRAND_PRIMARY,
                    subtitle_color=BRAND_WHITE, notes=""):
    """Add a branded title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_color)

    add_textbox(slide, 0.5, 1.2, 9.0, 1.2, title,
                font_size=54, bold=True, color=title_color,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 2.5, 9.0, 0.6, subtitle,
                font_size=24, color=subtitle_color,
                alignment=PP_ALIGN.CENTER)

    if extra_line:
        add_textbox(slide, 0.5, 4.6, 9.0, 0.5, extra_line,
                    font_size=14, color=BRAND_GRAY,
                    alignment=PP_ALIGN.CENTER)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_section_header_slide(prs, section_title, section_number="",
                             bg_color=BRAND_PRIMARY):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_color)

    if section_number:
        add_textbox(slide, 0.5, 1.5, 9.0, 0.6, section_number,
                    font_size=20, color=BRAND_WHITE,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.5, 2.1, 9.0, 1.0, section_title,
                font_size=44, bold=True, color=BRAND_WHITE,
                alignment=PP_ALIGN.CENTER)

    return slide


def add_thank_you_slide(prs, contact_name="[Your Name]",
                        contact_email="[email@studio.com]",
                        contact_url="[infinitevoyage.com]"):
    """Add a closing thank-you slide with contact info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_DARK)

    add_textbox(slide, 0.5, 1.2, 9.0, 1.0, "Thank You",
                font_size=54, bold=True, color=BRAND_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    contact = f"Questions?\n\n{contact_name}\n{contact_email}\n{contact_url}"
    txbox = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.8), Inches(7.0), Inches(2.0)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = contact
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.color.rgb = BRAND_WHITE
            run.font.size = Pt(16)

    slide.notes_slide.notes_text_frame.text = (
        "Thank the audience. Open the floor for Q&A. "
        "Have backup slides ready for detailed questions."
    )
    return slide


# ===========================================================================
# 1. Publisher Pitch Template (8 slides)
# ===========================================================================


def generate_publisher_pitch(output_dir):
    """Generate publisher-pitch-template.pptx -- 8 slides focused on gameplay and market."""
    prs = new_presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "INFINITE VOYAGE",
        "Publisher Pitch Deck",
        f"[Studio Name] | {TODAY_DISPLAY}",
        notes=(
            "Welcome the audience. Introduce yourself and your studio briefly. "
            "This is a 5-minute pitch focused on gameplay and market opportunity."
        ),
    )

    # Slide 2: The Hook
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_DARK)
    add_textbox(slide, 0.5, 1.0, 9.0, 0.5, "THE HOOK",
                font_size=18, bold=True, color=BRAND_SECONDARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(
        slide, 0.8, 1.8, 8.4, 1.5,
        '"What if the greatest discovery in the galaxy\n'
        'was also its most dangerous secret?"',
        font_size=32, bold=True, color=BRAND_WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, 1.0, 3.8, 8.0, 0.8,
        "[Replace with your game's most compelling one-liner.\n"
        "Use concept art or gameplay screenshot as background.]",
        font_size=14, color=BRAND_GRAY, alignment=PP_ALIGN.CENTER,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Open strong. Deliver the hook line with conviction. Pause for effect. "
        "This should make the audience lean forward. Show your most impressive visual."
    )

    # Slide 3: Gameplay
    add_bullet_slide(
        prs,
        "Core Gameplay",
        [
            "Genre: [e.g., Narrative Space Exploration RPG]",
            "Core Loop: Explore > Discover > Research > Upgrade > Explore Further",
            "Key Mechanic: [Your unique selling mechanic]",
            "Session Length: [e.g., 30-60 min per session, 40-60 hrs total]",
            "Comparable Titles: [e.g., Outer Wilds meets Mass Effect]",
        ],
        speaker_notes=(
            "Walk through the core gameplay loop. Show a gameplay video or animated "
            "mockup if possible. Emphasize what makes your game DIFFERENT from competitors. "
            "Spend 1-2 minutes here -- this is the heart of the pitch."
        ),
    )

    # Slide 4: Market Opportunity
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Market Opportunity",
                font_size=36, bold=True, color=BRAND_DARK)

    # Market data table
    table_shape = slide.shapes.add_table(4, 3, Inches(0.8), Inches(1.2),
                                         Inches(8.4), Inches(2.5))
    table = table_shape.table

    headers = ["Market Segment", "Size", "Growth"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    market_data = [
        ["[Primary Genre Market]", "[$X B]", "[+X% YoY]"],
        ["[Secondary Genre Market]", "[$X B]", "[+X% YoY]"],
        ["[Your Niche]", "[$X M]", "[+X% YoY]"],
    ]
    for r, row_data in enumerate(market_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    add_textbox(slide, 0.8, 4.0, 8.4, 0.8,
                "[Add TAM/SAM/SOM analysis. Reference comparable title sales.]",
                font_size=12, color=BRAND_GRAY)
    slide.notes_slide.notes_text_frame.text = (
        "Present market data credibly. Reference actual sales figures from comparable "
        "titles (SteamSpy, NPD, industry reports). Show your game fits a growing market."
    )

    # Slide 5: Monetization
    add_bullet_slide(
        prs,
        "Monetization Strategy",
        [
            "Business Model: [Premium / F2P / Premium + DLC]",
            "Price Point: [$XX.XX]",
            "DLC / Expansion Plans: [Post-launch content roadmap]",
            "Revenue Projection: [$X M first year based on comparable titles]",
            "Platform Split: [PC X%, Console X%, Other X%]",
        ],
        speaker_notes=(
            "Be specific about your revenue expectations and the data behind them. "
            "Publishers want to see realistic projections, not fantasies. "
            "Reference comparable title performance."
        ),
    )

    # Slide 6: Team
    add_bullet_slide(
        prs,
        "The Team",
        [
            "[Lead Name] - Game Director: [Credentials / shipped titles]",
            "[Lead Name] - Lead Programmer: [Credentials / tech expertise]",
            "[Lead Name] - Art Director: [Credentials / visual portfolio]",
            "[Lead Name] - Narrative Designer: [Credentials / published work]",
            "Team Size: [X core + X contractors]",
        ],
        speaker_notes=(
            "Highlight relevant experience. Publishers invest in teams, not just ideas. "
            "Mention any shipped titles, industry awards, or notable credentials."
        ),
    )

    # Slide 7: Timeline
    add_bullet_slide(
        prs,
        "Development Timeline",
        [
            "Pre-Production: [Dates] -- Completed",
            "Prototype / Vertical Slice: [Dates] -- [Status]",
            "Alpha: [Dates] -- [Status]",
            "Beta: [Dates] -- [Status]",
            "Launch: [Date] -- Target",
            "Post-Launch DLC: [Dates] -- Planned",
        ],
        speaker_notes=(
            "Show a realistic timeline. Highlight what is already done to demonstrate "
            "progress and reduce risk perception. Be honest about current status."
        ),
    )

    # Slide 8: The Ask
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_PRIMARY)
    add_textbox(slide, 0.5, 0.8, 9.0, 0.8, "The Ask",
                font_size=48, bold=True, color=BRAND_WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 1.8, 8.0, 0.8, "[Partnership Type: Publishing Deal]",
                font_size=28, bold=True, color=BRAND_SECONDARY,
                alignment=PP_ALIGN.CENTER)
    bullets_text = (
        "What We Need:\n"
        "  - [Funding: $X M for Y months of development]\n"
        "  - [Marketing support and platform relationships]\n"
        "  - [QA and localization]\n\n"
        "What You Get:\n"
        "  - [Revenue share: X% publisher / Y% developer]\n"
        "  - [Exclusive platform window: X months]\n"
        "  - [Sequel rights / first look deal]"
    )
    txbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7.0), Inches(2.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = bullets_text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BRAND_WHITE
            run.font.size = Pt(14)
    slide.notes_slide.notes_text_frame.text = (
        "Be specific about what you need and what you offer in return. "
        "Leave with a clear next step: 'We would love to schedule a follow-up "
        "to discuss partnership terms.'"
    )

    filepath = os.path.join(output_dir, "publisher-pitch-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


# ===========================================================================
# 2. Investor Pitch Template (12 slides)
# ===========================================================================


def generate_investor_pitch(output_dir):
    """Generate investor-pitch-template.pptx -- 12 slides focused on market and financials."""
    prs = new_presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "INFINITE VOYAGE",
        "Investor Presentation",
        f"[Studio Name] | {TODAY_DISPLAY} | Confidential",
        notes="Welcome investors. Set the stage: brief studio intro, then into the pitch.",
    )

    # Slide 2: Problem
    add_bullet_slide(
        prs,
        "The Problem",
        [
            "[Identify the gap in the market]",
            "[What are players not getting from current games?]",
            "[What frustration or unmet desire exists?]",
            "[Why is now the right time to solve this?]",
        ],
        speaker_notes=(
            "Frame the problem in terms investors understand: underserved market, "
            "growing demand, shift in player behavior. Use data to support claims."
        ),
    )

    # Slide 3: Solution
    add_bullet_slide(
        prs,
        "Our Solution",
        [
            "[How does your game solve the problem?]",
            "[What is the unique value proposition?]",
            "[Why will players choose this over alternatives?]",
            "[What is the emotional hook?]",
        ],
        speaker_notes=(
            "Describe the game as the answer to the problem. Keep it high-level. "
            "Save gameplay details for the next slide."
        ),
    )

    # Slide 4: Market Size
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Market Size",
                font_size=36, bold=True, color=BRAND_DARK)

    # TAM / SAM / SOM
    metrics = [
        ("TAM", "[Total Addressable Market: $XX B]", BRAND_PRIMARY),
        ("SAM", "[Serviceable Addressable Market: $XX B]", BRAND_SECONDARY),
        ("SOM", "[Serviceable Obtainable Market: $XX M]", BRAND_GREEN),
    ]
    y = 1.3
    for label, description, color in metrics:
        add_textbox(slide, 1.0, y, 1.5, 0.5, label,
                    font_size=28, bold=True, color=color)
        add_textbox(slide, 2.6, y + 0.05, 6.5, 0.4, description,
                    font_size=16, color=BRAND_DARK)
        y += 0.7

    add_textbox(slide, 0.8, 3.8, 8.4, 1.0,
                "[Add supporting data: industry reports, growth trends, "
                "comparable title revenues. Cite sources.]",
                font_size=12, color=BRAND_GRAY)
    slide.notes_slide.notes_text_frame.text = (
        "Investors want to see market size and your realistic capture rate. "
        "Use credible sources (Newzoo, SuperData, SteamSpy). Be conservative."
    )

    # Slide 5: Product
    add_bullet_slide(
        prs,
        "The Product",
        [
            "Genre: [e.g., Narrative Space Exploration RPG]",
            "Platform: [e.g., PC, PS5, Xbox Series X/S, Switch]",
            "Core Loop: [Brief description of gameplay loop]",
            "Unique Selling Point: [What makes this game special?]",
            "Current Status: [Prototype / Alpha / Beta]",
        ],
        speaker_notes=(
            "Show the game. Use gameplay footage, screenshots, or a trailer. "
            "Investors need to SEE what they are funding."
        ),
    )

    # Slide 6: Traction
    add_bullet_slide(
        prs,
        "Traction & Validation",
        [
            "Playtest Results: [X players tested, Y% positive feedback]",
            "Wishlists: [X Steam wishlists / social media followers]",
            "Awards / Recognition: [Festival selections, press coverage]",
            "Community: [Discord members, subreddit subscribers]",
            "Partnerships: [Any existing deals or LOIs]",
        ],
        speaker_notes=(
            "Traction reduces risk for investors. Show real numbers. "
            "Even small numbers are better than none -- they prove validation."
        ),
    )

    # Slide 7: Business Model
    add_bullet_slide(
        prs,
        "Business Model",
        [
            "Revenue Model: [Premium / F2P + MTX / Subscription]",
            "Price Point: [$XX.XX at launch]",
            "DLC Strategy: [X expansions planned at $XX each]",
            "Platform Revenue Split: [Standard 70/30 or negotiated]",
            "Merchandise / Licensing: [Additional revenue streams]",
        ],
        speaker_notes=(
            "Be clear about how the game makes money. Show you have thought "
            "beyond launch day. Multiple revenue streams reduce investor risk."
        ),
    )

    # Slide 8: Competition
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Competitive Landscape",
                font_size=36, bold=True, color=BRAND_DARK)

    # Comparison table
    table_shape = slide.shapes.add_table(5, 4, Inches(0.6), Inches(1.2),
                                         Inches(8.8), Inches(3.0))
    table = table_shape.table
    comp_headers = ["Feature", "Our Game", "[Competitor A]", "[Competitor B]"]
    for i, h in enumerate(comp_headers):
        table.cell(0, i).text = h
    comp_data = [
        ["[Feature 1]", "[Advantage]", "[Weakness]", "[Weakness]"],
        ["[Feature 2]", "[Advantage]", "[Neutral]", "[Advantage]"],
        ["[Feature 3]", "[Advantage]", "[Weakness]", "[Neutral]"],
        ["Price", "[$XX]", "[$XX]", "[$XX]"],
    ]
    for r, row_data in enumerate(comp_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    slide.notes_slide.notes_text_frame.text = (
        "Show you know the competition. Highlight your advantages honestly. "
        "Investors will research competitors -- be prepared for follow-up questions."
    )

    # Slide 9: Team
    add_bullet_slide(
        prs,
        "The Team",
        [
            "[CEO/Director]: [X years experience, shipped titles]",
            "[CTO/Lead Dev]: [Technical expertise, previous projects]",
            "[Creative Lead]: [Portfolio, awards, publications]",
            "[Business/Marketing]: [Industry connections, launch experience]",
            "Advisory Board: [Notable advisors if applicable]",
        ],
        speaker_notes=(
            "Investors fund teams. Highlight relevant experience, shipped titles, "
            "and why THIS team can execute THIS vision."
        ),
    )

    # Slide 10: Financials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Financial Projections",
                font_size=36, bold=True, color=BRAND_DARK)

    # Financial table
    fin_table_shape = slide.shapes.add_table(5, 4, Inches(0.6), Inches(1.2),
                                              Inches(8.8), Inches(2.8))
    fin_table = fin_table_shape.table
    fin_headers = ["", "Year 1", "Year 2", "Year 3"]
    for i, h in enumerate(fin_headers):
        fin_table.cell(0, i).text = h
    fin_data = [
        ["Revenue", "[$X M]", "[$X M]", "[$X M]"],
        ["Costs", "[$X M]", "[$X M]", "[$X M]"],
        ["Net Profit", "[$X M]", "[$X M]", "[$X M]"],
        ["ROI", "[X%]", "[X%]", "[X%]"],
    ]
    for r, row_data in enumerate(fin_data, 1):
        for c, val in enumerate(row_data):
            fin_table.cell(r, c).text = val

    add_textbox(slide, 0.6, 4.2, 8.8, 0.6,
                "[Include assumptions: unit sales, conversion rates, platform splits. "
                "Be conservative and cite comparable title performance.]",
                font_size=11, color=BRAND_GRAY)
    slide.notes_slide.notes_text_frame.text = (
        "Be realistic. Over-promising kills credibility. Show conservative, "
        "base, and optimistic scenarios. Explain your assumptions."
    )

    # Slide 11: The Ask
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_PRIMARY)
    add_textbox(slide, 0.5, 0.6, 9.0, 0.8, "The Ask",
                font_size=48, bold=True, color=BRAND_WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.0, 1.6, 8.0, 0.8, "[$X.X M] Seed / Series A",
                font_size=36, bold=True, color=BRAND_SECONDARY,
                alignment=PP_ALIGN.CENTER)

    use_of_funds = (
        "Use of Funds:\n"
        "  - Development (X%): Engineering, art, audio\n"
        "  - Marketing (X%): Community, PR, launch campaign\n"
        "  - Operations (X%): Studio overhead, tools, QA\n"
        "  - Reserve (X%): Contingency buffer\n\n"
        "Terms: [Equity / Convertible Note / Revenue Share]\n"
        "Runway: [X months to launch + X months post-launch]"
    )
    txbox = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(7.0), Inches(2.8))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.text = use_of_funds
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BRAND_WHITE
            run.font.size = Pt(14)
    slide.notes_slide.notes_text_frame.text = (
        "Be specific. Show exactly where the money goes. Have a detailed budget "
        "breakdown ready in backup slides. State your terms clearly."
    )

    # Slide 12: Thank You
    add_thank_you_slide(prs)

    filepath = os.path.join(output_dir, "investor-pitch-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


# ===========================================================================
# 3. Team Presentation Template (15 slides)
# ===========================================================================


def generate_team_presentation(output_dir):
    """Generate team-presentation-template.pptx -- 15 slides for internal deep-dive."""
    prs = new_presentation()

    # Slide 1: Title
    add_title_slide(
        prs,
        "INFINITE VOYAGE",
        "Team Presentation",
        f"[Studio Name] | {TODAY_DISPLAY} | Internal",
        notes="Welcome the team. This is a deep-dive into our current game vision and roadmap.",
    )

    # Slide 2: Agenda
    add_bullet_slide(
        prs,
        "Agenda",
        [
            "Vision & Pillars",
            "Core Gameplay Deep-Dive",
            "Systems Design",
            "Art Direction",
            "Narrative Overview",
            "Progression & Economy",
            "Technical Architecture",
            "Challenges & Solutions",
            "Roadmap & Milestones",
            "Action Items & Next Steps",
        ],
        speaker_notes="Walk through the agenda so everyone knows what to expect. ~20 minute presentation.",
    )

    # Slide 3: Vision
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_DARK)
    add_textbox(slide, 0.5, 0.8, 9.0, 0.5, "OUR VISION",
                font_size=18, bold=True, color=BRAND_SECONDARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(
        slide, 0.8, 1.6, 8.4, 1.5,
        "[Insert your vision statement here.\n"
        "What experience are we creating for players?\n"
        "What emotion should players feel?]",
        font_size=26, bold=True, color=BRAND_WHITE,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(slide, 1.0, 3.8, 8.0, 0.6,
                "[This should be the North Star that guides all design decisions.]",
                font_size=14, color=BRAND_GRAY, alignment=PP_ALIGN.CENTER)
    slide.notes_slide.notes_text_frame.text = (
        "Reaffirm the vision. Every team member should be able to recite this. "
        "It guides every design, art, and technical decision."
    )

    # Slide 4: Design Pillars
    add_bullet_slide(
        prs,
        "Design Pillars",
        [
            "Pillar 1: [e.g., Exploration -- Vast, hand-crafted worlds to discover]",
            "Pillar 2: [e.g., Discovery -- Ancient secrets with meaningful reveals]",
            "Pillar 3: [e.g., Consequence -- Player choices reshape the galaxy]",
            "Pillar 4: [e.g., Craftsmanship -- Every system polished, no filler]",
            "",
            "[Each pillar should be testable: 'Does this feature serve a pillar?']",
        ],
        speaker_notes=(
            "Review each pillar. Give examples of how they guide decisions. "
            "If a feature does not serve a pillar, we should question whether it belongs."
        ),
    )

    # Slide 5: Core Gameplay Loop
    add_bullet_slide(
        prs,
        "Core Gameplay Loop",
        [
            "EXPLORE: Navigate star systems, land on planets, scan for points of interest",
            "DISCOVER: Investigate ruins, solve environmental puzzles, find artifacts",
            "RESEARCH: Analyze discoveries, unlock technology, piece together history",
            "UPGRADE: Improve ship systems, unlock new abilities, expand reach",
            "RETURN: New upgrades open previously inaccessible areas",
            "",
            "[Session time target: 30-60 minutes per loop cycle]",
        ],
        speaker_notes=(
            "Walk through each phase of the loop in detail. Explain how they connect. "
            "This is the minute-to-minute experience players will have."
        ),
    )

    # Slide 6: Combat System
    add_bullet_slide(
        prs,
        "Combat & Interaction",
        [
            "Combat Style: [e.g., Turn-based tactical with action economy]",
            "Actions Per Turn: [e.g., 2 actions + 1 bonus action]",
            "Key Mechanics: [Positioning, elemental weakness, combo system]",
            "Difficulty Scaling: [Adaptive or fixed difficulty levels]",
            "Non-Combat Resolution: [Dialogue, stealth, hacking alternatives]",
        ],
        speaker_notes=(
            "Detail the combat system. Explain the design intent behind each mechanic. "
            "Address how non-combat players can still progress."
        ),
    )

    # Slide 7: Systems Design
    add_bullet_slide(
        prs,
        "Systems Overview",
        [
            "Inventory: [Weight-based / Slot-based / Grid-based]",
            "Crafting: [Material gathering + recipe discovery + upgrade tiers]",
            "Skill Trees: [Per-class trees with branching specializations]",
            "Equipment: [Slot system with rarity tiers and set bonuses]",
            "Economy: [Gold + crafting materials + research tokens]",
            "AI: [State machine enemies, companion behavior system]",
        ],
        speaker_notes=(
            "High-level overview of all systems. Each system has its own detailed "
            "design doc. Focus on how they interconnect."
        ),
    )

    # Slide 8: Art Direction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Art Direction",
                font_size=36, bold=True, color=BRAND_DARK)
    art_points = [
        "Visual Style: [e.g., Painterly sci-fi with bioluminescent accents]",
        "Color Palette: [Primary, secondary, accent colors defined]",
        "Character Design: [Silhouette language, faction color coding]",
        "Environment: [Biome types, architectural styles per civilization]",
        "UI: [Minimalist HUD, diegetic where possible]",
    ]
    y = 1.2
    for point in art_points:
        add_textbox(slide, 0.9, y, 8.2, 0.4, f"  {point}",
                    font_size=16, color=BRAND_DARK)
        y += 0.45

    add_textbox(slide, 0.6, 4.0, 8.8, 0.8,
                "[Add concept art, mood boards, and reference images to this slide.\n"
                "Replace this placeholder with visual examples.]",
                font_size=12, color=BRAND_GRAY, alignment=PP_ALIGN.CENTER)
    slide.notes_slide.notes_text_frame.text = (
        "Show concept art and references. The art direction should be visible, "
        "not just described. Include mood boards and style reference images."
    )

    # Slide 9: Narrative Overview
    add_bullet_slide(
        prs,
        "Narrative Overview",
        [
            "Setting: [Time, place, universe rules]",
            "Premise: [Central conflict and player motivation]",
            "Structure: [Three acts / Open-ended / Episodic]",
            "Characters: [Protagonist, antagonist, key NPCs]",
            "Branching: [Player choice impact on story and world]",
            "Tone: [Serious/humor ratio, darkness level, emotional targets]",
        ],
        speaker_notes=(
            "Summarize the story without spoilers for those who want to discover it fresh. "
            "Focus on structure and emotional beats."
        ),
    )

    # Slide 10: Progression & Economy
    add_bullet_slide(
        prs,
        "Progression & Economy",
        [
            "Level Cap: [e.g., Level 50]",
            "Progression Curve: [Linear / Exponential / S-curve]",
            "Currency: [Types and exchange rates]",
            "Faucets: [Enemy drops, quests, chests, salvage]",
            "Sinks: [Purchases, upgrades, repairs, cosmetics]",
            "Balance Target: [Slight positive flow at all levels]",
        ],
        speaker_notes=(
            "Reference the data-modeler spreadsheets for detailed numbers. "
            "Focus on design intent and balance philosophy."
        ),
    )

    # Slide 11: Technical Architecture
    add_bullet_slide(
        prs,
        "Technical Architecture",
        [
            "Engine: [e.g., Unreal Engine 5 / Unity / Custom]",
            "Target Platforms: [PC, PS5, Xbox Series X/S]",
            "Performance: [60 FPS target, 4K support]",
            "Networking: [Single-player with optional co-op]",
            "Tools: [Custom level editor, data pipeline, CI/CD]",
            "Key Technical Risks: [Streaming, AI performance, save system]",
        ],
        speaker_notes=(
            "Address technical feasibility. Highlight any prototype work already done. "
            "Be honest about technical risks and mitigation strategies."
        ),
    )

    # Slide 12: Challenges & Solutions
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, BRAND_WHITE)
    add_textbox(slide, 0.6, 0.3, 8.8, 0.7, "Challenges & Solutions",
                font_size=36, bold=True, color=BRAND_DARK)

    # Two-column layout
    challenge_items = [
        ("Challenge", "Solution"),
        ("[Content scope creep]", "[Pillar-based feature filtering]"),
        ("[Balance complexity]", "[Automated balance testing via data-modeler]"),
        ("[Technical risk: streaming]", "[Prototype completed, proven at scale]"),
        ("[Team bandwidth]", "[Phased roadmap with clear priorities]"),
    ]

    table_shape = slide.shapes.add_table(
        len(challenge_items), 2, Inches(0.6), Inches(1.2), Inches(8.8), Inches(3.0)
    )
    table = table_shape.table
    for r, (challenge, solution) in enumerate(challenge_items):
        table.cell(r, 0).text = challenge
        table.cell(r, 1).text = solution

    slide.notes_slide.notes_text_frame.text = (
        "Be transparent about challenges. Showing you have identified risks and "
        "have mitigation plans builds team confidence."
    )

    # Slide 13: Roadmap
    add_bullet_slide(
        prs,
        "Roadmap & Milestones",
        [
            "Milestone 1 (Completed): [Pre-production -- concept, prototypes, GDD v1]",
            "Milestone 2 (Current): [Vertical slice -- playable demo of core loop]",
            "Milestone 3: [Alpha -- all core systems functional]",
            "Milestone 4: [Beta -- content complete, polish and bug fixing]",
            "Milestone 5: [Gold -- release candidate, certification]",
            "Milestone 6: [Launch + post-launch content plan]",
        ],
        speaker_notes=(
            "Show where we are on the roadmap. Be specific about dates and deliverables. "
            "Highlight what is already done to build momentum."
        ),
    )

    # Slide 14: Action Items
    add_bullet_slide(
        prs,
        "Action Items & Next Steps",
        [
            "[Action 1]: [Owner] -- Due [Date]",
            "[Action 2]: [Owner] -- Due [Date]",
            "[Action 3]: [Owner] -- Due [Date]",
            "[Action 4]: [Owner] -- Due [Date]",
            "[Action 5]: [Owner] -- Due [Date]",
            "",
            "Next Team Sync: [Date and time]",
        ],
        speaker_notes=(
            "Assign specific, measurable action items with owners and deadlines. "
            "Every meeting should end with clear next steps."
        ),
    )

    # Slide 15: Thank You / Q&A
    add_thank_you_slide(prs)

    filepath = os.path.join(output_dir, "team-presentation-template.pptx")
    prs.save(filepath)
    print(f"  Created: {filepath}")


# ===========================================================================
# CLI entry point
# ===========================================================================


def main():
    parser = argparse.ArgumentParser(
        description="Generate Infinite Voyage pitch deck presentation templates (.pptx).",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python generate_templates.py\n"
            "  python generate_templates.py --output-dir ./output\n"
        ),
    )
    parser.add_argument(
        "--output-dir",
        default=os.path.dirname(os.path.abspath(__file__)),
        help="Directory where .pptx files will be written (default: script directory).",
    )
    args = parser.parse_args()

    output_dir = os.path.abspath(args.output_dir)
    os.makedirs(output_dir, exist_ok=True)

    print(f"Generating pitch deck templates in: {output_dir}\n")

    generate_publisher_pitch(output_dir)
    generate_investor_pitch(output_dir)
    generate_team_presentation(output_dir)

    print("\nDone. All pitch deck templates generated successfully.")


if __name__ == "__main__":
    main()
